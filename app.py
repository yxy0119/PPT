import os
import uuid
import tempfile
import requests
import json
import base64
from io import BytesIO
from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Pt

app = Flask(__name__)


def get_file_from_request(param_name):
    """
    从请求中获取文件对象，优先从上传文件获取，其次从JSON中的URL下载
    返回 (file-like object, filename)
    """
    # 1. 尝试从上传的文件中获取（multipart/form-data）
    file = request.files.get(param_name)
    if file and file.filename:
        return file, file.filename

    # 2. 尝试从 JSON 中获取 URL
    if request.is_json:
        data = request.get_json()
        url = data.get(param_name) if isinstance(data, dict) else None
        if url:
            try:
                resp = requests.get(url, timeout=30)
                resp.raise_for_status()
                filename = url.split('/')[-1].split('?')[0] or 'file.pptx'
                return BytesIO(resp.content), filename
            except Exception as e:
                print(f"下载文件失败: {e}")
                return None, None
    return None, None


def prepare_file_response(file_path, download_name):
    """
    根据请求参数 output 决定返回 JSON 还是文件
    如果请求中有 output=json，则返回 base64 编码的 JSON
    否则直接返回文件
    """
    output_format = request.args.get('output') or request.form.get('output')
    if output_format == 'json':
        with open(file_path, 'rb') as f:
            content = f.read()
        os.unlink(file_path)
        return jsonify({
            'filename': download_name,
            'base64': base64.b64encode(content).decode('utf-8')
        })
    else:
        return send_file(file_path, as_attachment=True, download_name=download_name)


@app.route('/')
def index():
    return 'PPT Service is running!'


@app.route('/openapi.json', methods=['GET'])
def openapi_spec():
    spec = {
        "openapi": "3.0.0",
        "info": {"title": "PPT 工具", "version": "1.0.0"},
        "servers": [{"url": "https://ppt-production-cca7.up.railway.app"}],
        "paths": {
            "/extract_template": {
                "post": {
                    "operationId": "extract_template",
                    "parameters": [
                        {
                            "name": "output",
                            "in": "query",
                            "required": False,
                            "schema": {"type": "string"}
                        }
                    ],
                    "requestBody": {
                        "content": {
                            "multipart/form-data": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "file": {"type": "string", "format": "binary"}
                                    },
                                    "required": ["file"]
                                }
                            }
                        }
                    },
                    "responses": {
                        "200": {
                            "description": "OK",
                            "content": {
                                "application/json": {
                                    "schema": {
                                        "type": "object",
                                        "properties": {
                                            "filename": {"type": "string"},
                                            "base64": {"type": "string"}
                                        }
                                    }
                                }
                            }
                        }
                    }
                }
            },
            "/generate_ppt": {
                "post": {
                    "operationId": "generate_ppt",
                    "parameters": [
                        {
                            "name": "output",
                            "in": "query",
                            "required": False,
                            "schema": {"type": "string"}
                        }
                    ],
                    "requestBody": {
                        "content": {
                            "multipart/form-data": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "template": {"type": "string", "format": "binary"},
                                        "slides_data": {"type": "string"}
                                    },
                                    "required": ["template", "slides_data"]
                                }
                            }
                        }
                    },
                    "responses": {
                        "200": {
                            "description": "OK",
                            "content": {
                                    "application/json": {
                                    "schema": {
                                        "type": "object",
                                        "properties": {
                                            "filename": {"type": "string"},
                                            "base64": {"type": "string"}
                                        }
                                    }
                                }
                            }
                        }
                    }
                }
            }
        }
    }
    return jsonify(spec)


@app.route('/extract_template', methods=['POST'])
def extract_template():
    file_obj, filename = get_file_from_request('file')
    if not file_obj:
        return jsonify({'error': 'No file'}), 400

    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        tmp.write(file_obj.read())
        src_path = tmp.name

    try:
        prs = Presentation(src_path)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get('r:id')
            if rId is None:
                rId = prs.slides._sldIdLst[0].get(
                    '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        template_path = f'/tmp/{uuid.uuid4()}.pptx'
        prs.save(template_path)
    finally:
        os.unlink(src_path)

    return prepare_file_response(template_path, 'template.pptx')

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    template_obj, _ = get_file_from_request('template')
    if not template_obj:
        return jsonify({'error': 'Missing template file'}), 400

    if request.is_json:
        data = request.get_json()
        slides_data = data.get('slides_data', '[]')
    else:
        slides_data = request.form.get('slides_data', '[]')

    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        tmp.write(template_obj.read())
        template_path = tmp.name

    try:
        slides = json.loads(slides_data)
        prs = Presentation(template_path)
        layouts = prs.slide_layouts

        for slide_data in slides:
            layout_name = slide_data.get('layout', '标题和内容')
            matched_layout = None
            for layout in layouts:
                if layout_name in layout.name:
                    matched_layout = layout
                    break
            if not matched_layout:
                matched_layout = layouts[1]

            slide = prs.slides.add_slide(matched_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get('title', '')

            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    for point in slide_data.get('bullet_points', []):
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                    for fixed_text in slide_data.get('fixed_texts', []):
                        p = tf.add_paragraph()
                        p.text = fixed_text
                        p.font.size = Pt(10)
                    break

        output_path = f'/tmp/{uuid.uuid4()}.pptx'
        prs.save(output_path)
    finally:
        os.unlink(template_path)

    return prepare_file_response(output_path, 'generated.pptx')


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8080)
    
                             
